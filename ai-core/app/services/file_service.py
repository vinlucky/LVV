import os
import uuid
import shutil
import logging
from datetime import datetime, timezone
from pathlib import Path
from fastapi import UploadFile

from ..config import FILE_DIR, SKILLS_DIR
from ..database import save_file_record, get_files_by_conv, get_all_files, update_file_conv_id

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    "txt", "md", "doc", "docx", "xls", "xlsx",
    "pdf", "ppt", "pptx",
    "py", "js", "ts", "java", "c", "cpp", "go", "rs",
    "html", "htm", "tex", "rtf", "csv", "json",
    "yaml", "yml", "xml", "sql", "sh", "bat", "ps1",
    "mp3", "wav", "ogg", "flac", "aac", "m4a", "wma", "opus", "webm",
    "mp4", "avi", "mov", "mkv", "wmv", "flv",
}

AUDIO_EXTENSIONS = {"mp3", "wav", "ogg", "flac", "aac", "m4a", "wma", "opus", "webm"}
OCR_CACHE_DIR = FILE_DIR / "ocrtmp"
OCR_CACHE_DIR.mkdir(parents=True, exist_ok=True)
os.environ["PADDLEOCR_HOME"] = str(OCR_CACHE_DIR)
os.environ["PADDLEX_HOME"] = str(FILE_DIR / "paddlex_models")


def _get_conv_dir(conv_id: str, mode: str = "chat") -> Path:
    existing = list(FILE_DIR.glob(f"{mode}/*/*_{conv_id[:8]}"))
    if existing:
        return existing[0]
    now = datetime.now(timezone.utc)
    date_str = now.strftime("%Y-%m-%d_%H-%M-%S")
    dir_name = f"{date_str}_{conv_id[:8]}"
    conv_dir = FILE_DIR / mode / dir_name
    conv_dir.mkdir(parents=True, exist_ok=True)
    return conv_dir


def get_file_extension(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def normalize_filename(filename: str, conv_dir: Path) -> str:
    base = Path(filename).stem
    ext = get_file_extension(filename)
    counter = 0
    new_name = filename
    while (conv_dir / new_name).exists():
        counter += 1
        new_name = f"{base}_{counter}.{ext}"
    return new_name


async def save_upload_file(file: UploadFile, conv_id: str | None = None, mode: str = "chat") -> dict:
    conv_dir = _get_conv_dir(conv_id, mode) if conv_id else FILE_DIR / "unfiled"
    conv_dir.mkdir(parents=True, exist_ok=True)

    original_filename = file.filename or "unnamed"
    stored_filename = normalize_filename(original_filename, conv_dir)
    file_path = conv_dir / stored_filename

    content = await file.read()
    file_size = len(content)
    with open(file_path, "wb") as f:
        f.write(content)

    file_format = get_file_extension(original_filename)
    relative_path = str(file_path.relative_to(FILE_DIR.parent))

    record = save_file_record(
        filename=original_filename,
        stored_filename=stored_filename,
        file_path=str(file_path),
        file_format=file_format,
        file_size=file_size,
        conv_id=conv_id,
        mode=mode,
    )

    result = {
        "filename": original_filename,
        "stored_filename": stored_filename,
        "file_path": str(file_path),
        "relative_path": relative_path,
        "file_format": file_format,
        "file_size": file_size,
        "conv_id": conv_id,
        "download_url": f"/files/download/{relative_path.replace(chr(92), '/')}",
        "mode": mode,
    }

    logger.info(f"[file] Upload saved: '{original_filename}' ({file_size} bytes, mode={mode}, conv_id={conv_id})")

    return result


def save_generated_file(content: str, filename: str, file_format: str = "md",
                         conv_id: str | None = None, mode: str = "chat") -> dict:
    conv_dir = _get_conv_dir(conv_id, mode) if conv_id else FILE_DIR / "unfiled"
    conv_dir.mkdir(parents=True, exist_ok=True)

    if not filename.endswith(f".{file_format}"):
        filename = f"{filename}.{file_format}"
    stored_filename = normalize_filename(filename, conv_dir)
    file_path = conv_dir / stored_filename

    logger.info(f"[file] Generating file: '{filename}' (format={file_format}, mode={mode}, conv_id={conv_id})")

    if file_format == "docx":
        logger.info(f"[file] Writing Word document: {filename}")
        _write_docx(content, file_path)
        logger.info(f"[file] Word document written: {file_path} ({file_path.stat().st_size if file_path.exists() else 0} bytes)")
    elif file_format == "xlsx":
        logger.info(f"[file] Writing Excel spreadsheet: {filename}")
        _write_xlsx(content, file_path)
        logger.info(f"[file] Excel spreadsheet written: {file_path} ({file_path.stat().st_size if file_path.exists() else 0} bytes)")
    else:
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(content)
        logger.info(f"[file] Text file written: {file_path} ({file_path.stat().st_size if file_path.exists() else 0} bytes)")

    file_size = file_path.stat().st_size if file_path.exists() else 0
    relative_path = str(file_path.relative_to(FILE_DIR.parent))

    record = save_file_record(
        filename=filename,
        stored_filename=stored_filename,
        file_path=str(file_path),
        file_format=file_format,
        file_size=file_size,
        conv_id=conv_id,
        mode=mode,
        is_generated=True,
    )

    return {
        "filename": filename,
        "file_name": Path(filename).stem,
        "stored_filename": stored_filename,
        "file_path": str(file_path),
        "relative_path": relative_path,
        "file_format": file_format,
        "file_size": file_size,
        "conv_id": conv_id,
        "download_url": f"/files/download/{relative_path.replace(chr(92), '/')}",
        "mode": mode,
    }


def _write_docx(content: str, file_path: Path) -> None:
    try:
        from docx import Document
        doc = Document()
        for line in content.split("\n"):
            if line.startswith("# "):
                doc.add_heading(line[2:], level=1)
            elif line.startswith("## "):
                doc.add_heading(line[3:], level=2)
            elif line.startswith("### "):
                doc.add_heading(line[4:], level=3)
            elif line.strip():
                doc.add_paragraph(line)
        doc.save(str(file_path))
    except ImportError:
        logger.warning("python-docx not installed, saving as plain text with .docx extension")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(content)


def _write_xlsx(content: str, file_path: Path) -> None:
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        for row_idx, line in enumerate(content.split("\n"), 1):
            cells = line.split("\t")
            for col_idx, cell in enumerate(cells, 1):
                ws.cell(row=row_idx, column=col_idx, value=cell.strip())
        wb.save(str(file_path))
    except ImportError:
        logger.warning("openpyxl not installed, saving as plain text with .xlsx extension")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(content)


def read_text_file(file_path: str) -> str:
    path = Path(file_path)
    if not path.is_absolute():
        path = FILE_DIR.parent / file_path

    ext = path.suffix.lower()

    if ext in (".txt", ".md", ".py", ".js", ".ts", ".java", ".c", ".cpp", ".go", ".rs",
               ".html", ".htm", ".tex", ".rtf", ".csv", ".json", ".yaml", ".yml",
               ".xml", ".sql", ".sh", ".bat", ".ps1"):
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    if ext == ".pdf":
        return _read_pdf(path)

    if ext in (".docx", ".doc"):
        return _read_docx(path)

    if ext in (".xlsx", ".xls"):
        return _read_xlsx(path)

    if ext in (".pptx", ".ppt"):
        return _read_pptx(path)

    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


_paddleocr_instance = None
PADDLEOCR_MAX_PAGES = 100
PADDLEOCR_WORKERS = 2
OCR_LOW_CONFIDENCE_THRESHOLD = 0.5
VISION_OCR_MAX_PAGES = 30


def _ocr_page_worker(page_index: int, pdf_path: str, scale: float) -> tuple[int, str, float]:
    global _paddleocr_instance
    if _paddleocr_instance is None:
        from paddleocr import PaddleOCR
        logger.info(f"[paddleocr:worker] Initializing PaddleOCR (page {page_index + 1})...")
        _paddleocr_instance = PaddleOCR(use_angle_cls=False, ocr_version='PP-OCRv4')
        logger.info("[paddleocr:worker] PaddleOCR initialized")
    import numpy as np
    import pypdfium2 as pdfium
    pdf = pdfium.PdfDocument(pdf_path)
    page = pdf[page_index]
    bitmap = page.render(scale=scale)
    img = bitmap.to_pil()
    img_array = np.array(img)
    result = list(_paddleocr_instance.predict(img_array))
    if result:
        rec_texts = result[0].get("rec_texts", [])
        rec_scores = result[0].get("rec_scores", [])
        if rec_texts:
            avg_score = float(sum(rec_scores) / len(rec_scores)) if rec_scores else 0.0
            return (page_index, "\n".join(rec_texts), avg_score)
    return (page_index, "", 0.0)


def _get_paddleocr():
    global _paddleocr_instance
    if _paddleocr_instance is None:
        from paddleocr import PaddleOCR
        logger.info("[paddleocr] Initializing PaddleOCR (first call, downloading models if needed)...")
        _paddleocr_instance = PaddleOCR(use_angle_cls=False, ocr_version='PP-OCRv4')
        logger.info("[paddleocr] PaddleOCR initialized successfully")
    return _paddleocr_instance


def _paddleocr_sequential(path: Path, max_pages: int) -> tuple[list[str], list[float]]:
    import numpy as np
    import pypdfium2 as pdfium
    ocr = _get_paddleocr()
    pdf = pdfium.PdfDocument(str(path))
    text_parts = []
    page_scores = []
    for i in range(max_pages):
        page = pdf[i]
        bitmap = page.render(scale=1.5)
        img = bitmap.to_pil()
        img_array = np.array(img)
        result = list(ocr.predict(img_array))
        avg_score = 0.0
        if result:
            rec_texts = result[0].get("rec_texts", [])
            rec_scores = result[0].get("rec_scores", [])
            if rec_texts:
                text_parts.append("\n".join(rec_texts))
                avg_score = float(sum(rec_scores) / len(rec_scores)) if rec_scores else 0.0
        page_scores.append(avg_score)
        logger.info(f"[paddleocr] Page {i + 1}/{max_pages} done (confidence: {avg_score:.2f})")
    return text_parts, page_scores


def _read_pdf_with_paddleocr(path: Path) -> str:
    try:
        import numpy as np
        import pypdfium2 as pdfium
    except ImportError:
        logger.warning("[paddleocr] Dependencies not installed (paddleocr, pypdfium2, numpy)")
        return ""

    try:
        pdf = pdfium.PdfDocument(str(path))
        total_pages = len(pdf)
        max_pages = min(total_pages, PADDLEOCR_MAX_PAGES)
        pdf.close()

        if total_pages > PADDLEOCR_MAX_PAGES:
            logger.warning(f"[paddleocr] PDF has {total_pages} pages, only processing first {PADDLEOCR_MAX_PAGES}")

        logger.info(f"[paddleocr] Starting OCR for: {path.name} ({max_pages} pages, {PADDLEOCR_WORKERS} workers)")

        text_parts: list[str] = []
        page_scores: list[float] = []

        if max_pages > 1:
            try:
                from concurrent.futures import ProcessPoolExecutor, as_completed
                results: dict[int, tuple[str, float]] = {}
                with ProcessPoolExecutor(max_workers=PADDLEOCR_WORKERS) as executor:
                    futures = {
                        executor.submit(_ocr_page_worker, i, str(path), 1.5): i
                        for i in range(max_pages)
                    }
                    for future in as_completed(futures):
                        page_idx, text, score = future.result()
                        results[page_idx] = (text, score)
                        logger.info(f"[paddleocr] Page {page_idx + 1}/{max_pages} done ({len(text)} chars, confidence: {score:.2f})")
                for i in range(max_pages):
                    if i in results and results[i][0]:
                        text_parts.append(results[i][0])
                        page_scores.append(results[i][1])
            except Exception as e:
                logger.warning(f"[paddleocr] Parallel processing failed: {e}, falling back to sequential")
                text_parts, page_scores = _paddleocr_sequential(path, max_pages)
        else:
            text_parts, page_scores = _paddleocr_sequential(path, max_pages)

        if text_parts:
            overall_confidence = sum(page_scores) / len(page_scores) if page_scores else 0.0
            total_chars = sum(len(t) for t in text_parts)
            logger.info(f"[paddleocr] Completed: {len(text_parts)}/{max_pages} pages, {total_chars} chars, avg confidence: {overall_confidence:.2f}")

            if overall_confidence < OCR_LOW_CONFIDENCE_THRESHOLD:
                logger.warning(f"[paddleocr] Low confidence ({overall_confidence:.2f} < {OCR_LOW_CONFIDENCE_THRESHOLD}), likely unrecognized language or poor scan quality, will fall back to Vision LLM OCR")
                return ""

            return "\n\n".join(text_parts)
        else:
            logger.warning(f"[paddleocr] No text extracted from any page of {path.name}")
    except Exception as e:
        logger.warning(f"[paddleocr] Failed: {e}")

    return ""


def _read_pdf_with_vision_llm(path: Path) -> str:
    try:
        import base64
        from io import BytesIO
        import pypdfium2 as pdfium
        from PIL import Image
        from .llm_service import vision_ocr_with_fallback
    except ImportError as e:
        logger.warning(f"[vision_ocr] Dependencies not available: {e}")
        return ""

    try:
        pdf = pdfium.PdfDocument(str(path))
        total_pages = len(pdf)
        max_pages = min(total_pages, VISION_OCR_MAX_PAGES)

        if total_pages > VISION_OCR_MAX_PAGES:
            logger.warning(f"[vision_ocr] PDF has {total_pages} pages, only processing first {VISION_OCR_MAX_PAGES}")

        logger.info(f"[vision_ocr] Starting Vision LLM OCR for: {path.name} ({total_pages} pages, processing {max_pages})")
        text_parts = []
        for i in range(max_pages):
            page = pdf[i]
            bitmap = page.render(scale=2.0)
            img = bitmap.to_pil()

            max_dim = 1600
            width, height = img.size
            if max(width, height) > max_dim:
                scale_factor = max_dim / max(width, height)
                new_width = int(width * scale_factor)
                new_height = int(height * scale_factor)
                img = img.resize((new_width, new_height), Image.LANCZOS)

            buffer = BytesIO()
            img.save(buffer, format="JPEG", quality=85)
            img_b64 = base64.b64encode(buffer.getvalue()).decode("utf-8")

            logger.info(f"[vision_ocr] Processing page {i + 1}/{max_pages}...")
            page_text = vision_ocr_with_fallback(img_b64, page_num=i + 1, total_pages=max_pages)
            if page_text:
                text_parts.append(page_text)
                logger.info(f"[vision_ocr] Page {i + 1}/{max_pages} done ({len(page_text)} chars)")
            else:
                logger.info(f"[vision_ocr] Page {i + 1}/{max_pages} done (no text returned)")

        if text_parts:
            total_chars = sum(len(t) for t in text_parts)
            logger.info(f"[vision_ocr] Completed: {len(text_parts)}/{max_pages} pages extracted, {total_chars} chars total")
            return "\n\n".join(text_parts)
        else:
            logger.warning(f"[vision_ocr] No text extracted from any page of {path.name}")
    except Exception as e:
        logger.warning(f"[vision_ocr] Failed: {e}")

    return ""


def _read_pdf(path: Path) -> str:
    logger.info(f"[PDF] Reading file: {path.name} (size: {path.stat().st_size} bytes)")
    text = _read_pdf_with_paddleocr(path)
    if text:
        logger.info(f"[PDF] Successfully read via PaddleOCR: {path.name} ({len(text)} chars)")
        return text

    logger.info(f"[PDF] PaddleOCR returned empty, falling back to Vision LLM OCR: {path.name}")
    text = _read_pdf_with_vision_llm(path)
    if text:
        logger.info(f"[PDF] Successfully read via Vision LLM OCR: {path.name} ({len(text)} chars)")
        return text

    logger.error(f"[PDF] All OCR methods failed for: {path.name}")
    return ""


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except Exception as e:
        logger.warning(f"docx read failed: {e}")
        return f"[无法解析 DOCX 文件: {path.name}]"


def _read_xlsx(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"=== Sheet: {sheet_name} ===")
            rows_text = []
            for row in ws.iter_rows(values_only=True):
                row_str = "\t".join(str(c) if c is not None else "" for c in row)
                if row_str.strip():
                    rows_text.append(row_str)
            parts.append("\n".join(rows_text))
        return "\n\n".join(parts)
    except Exception as e:
        logger.warning(f"xlsx read failed: {e}")
        return f"[无法解析 XLSX 文件: {path.name}]"


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides):
            slide_text = [f"=== 幻灯片 {i + 1} ==="]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_text.append(t)
            parts.append("\n".join(slide_text))
        return "\n\n".join(parts)
    except Exception as e:
        logger.warning(f"pptx read failed: {e}")
        return f"[无法解析 PPTX 文件: {path.name}]"


def get_download_path(file_path: str) -> Path:
    path = Path(file_path)
    if not path.is_absolute():
        path = FILE_DIR.parent / file_path
    return path


def delete_file(file_path: str) -> bool:
    path = Path(file_path)
    if not path.is_absolute():
        path = FILE_DIR.parent / file_path
    if path.exists() and path.is_file():
        path.unlink()
        return True
    return False


VIDEO_EXTENSIONS = {"mp4", "avi", "mov", "mkv", "wmv", "flv", "webm", "m4v", "3gp"}


def extract_audio_from_video(video_path: str) -> str | None:
    import time as _time
    path = Path(video_path)
    if not path.exists():
        logger.error(f"[video] File not found: {video_path}")
        return None

    ext = path.suffix.lower()
    if ext.lstrip('.') not in VIDEO_EXTENSIONS:
        logger.info(f"[video] Not a video file: {ext}, skipping audio extraction")
        return None

    video_size_mb = path.stat().st_size / 1024 / 1024
    logger.info(f"[video] Starting audio extraction: {path.name} ({video_size_mb:.1f}MB, format={ext})")

    audio_output = path.with_suffix('.mp3')

    start_time = _time.time()
    try:
        from moviepy import VideoFileClip
        logger.info(f"[video] Loading video with MoviePy+imageio-ffmpeg: {path.name}")
        clip = VideoFileClip(str(path))
        video_duration = clip.duration
        logger.info(f"[video] Video duration: {video_duration:.1f}s ({video_duration/60:.1f}min)")

        if clip.audio is None:
            logger.warning(f"[video] No audio track found in: {path.name}")
            clip.close()
            return None

        logger.info(f"[video] Extracting audio track → {audio_output.name}")
        clip.audio.write_audiofile(str(audio_output), logger=None)
        clip.close()

        elapsed = _time.time() - start_time
        if audio_output.exists() and audio_output.stat().st_size > 0:
            audio_size_mb = audio_output.stat().st_size / 1024 / 1024
            logger.info(f"[video] Audio extracted successfully: {audio_output.name} ({audio_size_mb:.1f}MB) in {elapsed:.1f}s")
            return str(audio_output)
        else:
            logger.warning(f"[video] Audio extraction produced empty file: {audio_output} ({elapsed:.1f}s)")
            return None
    except ImportError:
        logger.warning("[video] moviepy not installed, falling back to ffmpeg directly")
        return _extract_audio_ffmpeg(video_path)
    except Exception as e:
        elapsed = _time.time() - start_time
        logger.warning(f"[video] MoviePy extraction failed after {elapsed:.1f}s: {e}, falling back to ffmpeg")
        return _extract_audio_ffmpeg(video_path)


def _extract_audio_ffmpeg(video_path: str) -> str | None:
    import subprocess
    import time as _time
    path = Path(video_path)
    audio_output = path.with_suffix('.mp3')

    try:
        import imageio_ffmpeg
        ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
        logger.info(f"[video] Using imageio-ffmpeg: {ffmpeg_exe}")
    except ImportError:
        ffmpeg_exe = "ffmpeg"
        logger.info("[video] imageio-ffmpeg not available, using system ffmpeg")

    try:
        video_size_mb = path.stat().st_size / 1024 / 1024
        logger.info(f"[video] Extracting audio via ffmpeg: {path.name} ({video_size_mb:.1f}MB)")
        start_time = _time.time()
        result = subprocess.run(
            [ffmpeg_exe, "-i", str(path), "-vn", "-acodec", "libmp3lame",
             "-q:a", "4", "-y", str(audio_output)],
            capture_output=True, text=True, timeout=300,
        )
        elapsed = _time.time() - start_time

        if result.returncode != 0:
            logger.error(f"[video] ffmpeg failed (exit={result.returncode}) after {elapsed:.1f}s: {result.stderr[:500]}")
            return None

        if audio_output.exists() and audio_output.stat().st_size > 0:
            audio_size_mb = audio_output.stat().st_size / 1024 / 1024
            logger.info(f"[video] ffmpeg audio extracted: {audio_output.name} ({audio_size_mb:.1f}MB) in {elapsed:.1f}s")
            return str(audio_output)
        else:
            logger.warning(f"[video] ffmpeg produced empty file after {elapsed:.1f}s")
            return None
    except FileNotFoundError:
        logger.error("[video] ffmpeg not found on system")
        return None
    except subprocess.TimeoutExpired:
        logger.error("[video] ffmpeg timed out after 300s")
        return None
    except Exception as e:
        logger.error(f"[video] ffmpeg extraction failed: {e}")
        return None